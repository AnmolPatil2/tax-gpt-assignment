from __future__ import annotations

import logging
import struct
from pathlib import Path

from app.retrieval import vector_store

logger = logging.getLogger(__name__)


def _extract_text_pptx(file_path: Path) -> list[tuple[int, str]]:
    """Extract text from .pptx files using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append((i + 1, "\n".join(texts)))
    return slides


def _extract_text_ole(file_path: Path) -> list[tuple[int, str]]:
    """Extract text from old .ppt (OLE2) files using olefile."""
    import olefile

    ole = olefile.OleFileIO(str(file_path))
    text_chunks = []

    for stream_name in ole.listdir():
        stream_path = "/".join(stream_name)
        if "PowerPoint Document" in stream_path or "Current User" in stream_path:
            try:
                data = ole.openstream(stream_name).read()
                extracted = _extract_text_from_ppt_binary(data)
                if extracted:
                    text_chunks.extend(extracted)
            except Exception:
                continue

    ole.close()

    if not text_chunks:
        return _extract_text_raw(file_path)

    # Group extracted text into pseudo-slides
    slides = []
    current_slide: list[str] = []
    slide_num = 1

    for chunk in text_chunks:
        chunk = chunk.strip()
        if not chunk:
            continue
        current_slide.append(chunk)
        # Heuristic: a title-like short chunk followed by content suggests a new slide
        if len(current_slide) > 5:
            slides.append((slide_num, "\n".join(current_slide)))
            current_slide = []
            slide_num += 1

    if current_slide:
        slides.append((slide_num, "\n".join(current_slide)))

    return slides


def _extract_text_from_ppt_binary(data: bytes) -> list[str]:
    """Parse PPT binary stream to find text runs (record type 4008 = TextBytesAtom,
    4000 = TextCharsAtom)."""
    texts = []
    offset = 0

    while offset < len(data) - 8:
        try:
            rec_ver_inst = struct.unpack_from("<H", data, offset)[0]
            rec_type = struct.unpack_from("<H", data, offset + 2)[0]
            rec_len = struct.unpack_from("<I", data, offset + 4)[0]
        except struct.error:
            break

        offset += 8

        if rec_type == 4008 and rec_len > 0 and rec_len < 100000:
            # TextBytesAtom: ASCII/Latin1 text
            try:
                text = data[offset : offset + rec_len].decode("latin-1", errors="ignore")
                text = text.strip()
                if text and len(text) > 1:
                    texts.append(text)
            except Exception:
                pass
        elif rec_type == 4000 and rec_len > 0 and rec_len < 200000:
            # TextCharsAtom: UTF-16LE text
            try:
                text = data[offset : offset + rec_len].decode("utf-16-le", errors="ignore")
                text = text.strip()
                if text and len(text) > 1:
                    texts.append(text)
            except Exception:
                pass

        if rec_len > 0:
            offset += rec_len
        else:
            offset += 1

    return texts


def _extract_text_raw(file_path: Path) -> list[tuple[int, str]]:
    """Last-resort: extract printable ASCII strings from the binary file."""
    with open(file_path, "rb") as f:
        data = f.read()

    import re

    strings = re.findall(rb"[\x20-\x7e]{10,}", data)
    decoded = [s.decode("ascii", errors="ignore").strip() for s in strings]
    # Filter out binary artifacts
    filtered = [s for s in decoded if len(s) > 15 and not s.startswith(("Microsoft", "Arial", "Times"))]

    if not filtered:
        return []

    combined = "\n".join(filtered)
    return [(1, combined)]


def ingest(file_path: str | Path) -> dict:
    file_path = Path(file_path)
    logger.info(f"Ingesting PPT: {file_path.name}")

    slides: list[tuple[int, str]] = []

    # Try python-pptx first (works for .pptx)
    try:
        slides = _extract_text_pptx(file_path)
        logger.info(f"Extracted {len(slides)} slides via python-pptx")
    except Exception as e:
        logger.info(f"python-pptx failed ({e}), trying OLE extraction")

    # Fallback to OLE parsing for old .ppt format
    if not slides:
        try:
            slides = _extract_text_ole(file_path)
            logger.info(f"Extracted {len(slides)} slide groups via OLE parsing")
        except Exception as e:
            logger.warning(f"OLE extraction failed ({e}), trying raw extraction")

    # Last resort: raw string extraction
    if not slides:
        slides = _extract_text_raw(file_path)
        logger.info(f"Extracted {len(slides)} text blocks via raw parsing")

    if not slides:
        logger.warning(f"Could not extract any text from {file_path.name}")
        return {"slides_processed": 0, "vector_chunks": 0}

    # Chunk long text blocks to stay within embedding model token limit
    from app.ingestion.pdf_loader import _chunk_text
    from app.config import settings

    ids = []
    texts = []
    metadatas = []
    chunk_idx = 0

    for slide_num, text in slides:
        if len(text) > 3000:
            chunks = _chunk_text(text, settings.chunk_size, settings.chunk_overlap)
        else:
            chunks = [text]
        for chunk in chunks:
            ids.append(f"ppt_{file_path.stem}_{chunk_idx}")
            texts.append(chunk)
            metadatas.append({
                "source_type": "ppt",
                "document": file_path.name,
                "slide": slide_num,
            })
            chunk_idx += 1

    vector_count = vector_store.add_documents(ids, texts, metadatas)
    logger.info(f"Added {vector_count} chunks to ChromaDB from {file_path.name}")

    return {
        "slides_processed": len(slides),
        "vector_chunks": vector_count,
    }
